#!/usr/bin/env python3
"""
deal_extract_helpers.py — file I/O primitives invoked by the /deal-sync
slash command.

The slash command is the workflow orchestrator and runs inside a Claude
Code session (subscription-backed). All AI work — reading new source
files, regenerating status/brief docs, regenerating the firm-context
pipeline section — happens in the LLM context window. This module does
NO API calls. It only handles Drive reads/writes, dedup state, and
filesystem moves.

Sub-commands (each prints structured output the slash command can parse):

    list-new-files <deal_id>
        Print JSON list of files in the deal's Drive folder modified
        since last_run, filtered to processable types, excluding
        _Ready/ contents and per-deal state files (status, brief,
        dashboard_entry).

    read-file <file_id>
        Print raw file contents (Drive media get + decoded text).

    read-deal-doc <deal_id> {status|brief}
        Print current status or brief doc contents.

    write-deal-doc <deal_id> {status|brief}
        Read new content from stdin, overwrite the corresponding doc
        in Drive (text/plain media update).

    write-actions-md <deal_id>
        Read dashboard_entry JSON from stdin, extract actions array,
        write a formatted 7-column markdown table to the local
        actions.md file (skips if file already has real content).
        Also syncs to Drive actions doc via Deal Sync Writer if
        ~/cos-pipeline-config-tomac/config/deal_sync.yaml exists.

    move-to-ready <file_id> <deal_folder_id>
        Move file into the deal's _Ready/ subfolder (creates folder
        if absent), mirroring the cos_transcript_hook _Ready/ pattern.

    mark-processed <deal_id> <file_id> {success|failed}
        Append/update the dedup record in deal_extract_state.json.

    update-last-run <deal_id>
        Set last_run for the deal to now (called after a successful
        cycle, even if no files were processed).

    list-deals [--filter <deal_id>,<deal_id>]
        Print JSON list of {deal_id, status_id, brief_id, folder_id}
        for each deal in deal_docs registry.

    regenerate-pipeline-section
        Read replacement section content from stdin, replace the body
        between AUTO-GENERATED-PIPELINE-START/END markers in the
        Firm Context doc.

Usage examples:
    python3 deal_extract_helpers.py list-deals
    python3 deal_extract_helpers.py list-new-files <deal_id>
    python3 deal_extract_helpers.py read-file 1abc...
    cat new_status.md | python3 deal_extract_helpers.py write-deal-doc <deal_id> status
"""
import argparse
import io
import json
import os
import pickle
import sys
from datetime import datetime
from pathlib import Path

import yaml
from googleapiclient.discovery import build
from googleapiclient.http import MediaInMemoryUpload, MediaIoBaseDownload

# ── Paths / config ────────────────────────────────────────────────────────────
CREDS_PATH = os.path.expanduser("~/credentials/gdrive_token.pickle")
DRIVE_DOCS_YAML = os.path.expanduser("~/dashboards/config/drive-docs.yaml")
EXTRACT_STATE_PATH = os.path.expanduser("~/dashboards/data/deal_extract_state.json")
EXTRACT_LOG_DIR = os.path.expanduser("~/dashboards/logs")
ERROR_LOG = os.path.join(EXTRACT_LOG_DIR, "extract_errors.log")
DEAL_SYNC_YAML = os.path.expanduser("~/cos-pipeline-config-tomac/config/deal_sync.yaml")

# Files in a deal folder we never feed back as inputs (they ARE the deal state)
EXCLUDED_NAME_SUFFIXES = (
    "_status.md",
    "_master_brief.md",
    "_dashboard_entry.json",
    " -- Status.md",
    " -- Master Brief.md",
    " -- Master Brief v2.md",
    " -- Dashboard Entry.json",
)

# Subfolder names we skip entirely when traversing recursively. _Ready/ holds
# already-processed source files; the others typically contain compiled outputs
# (decks, models) that aren't source intel.
EXCLUDED_SUBFOLDER_NAMES = {"_Ready", "Distributed", "Presentations"}

MAX_FOLDER_DEPTH = 5

# Mime types we accept as processable. Office formats are text-extracted via
# python-docx / openpyxl / python-pptx in drive_read_file_text().
MIME_DOCX = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
MIME_XLSX = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
MIME_PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
MIME_GDOC = "application/vnd.google-apps.document"

PROCESSABLE_MIMES = {
    "text/plain",
    "text/markdown",
    "application/pdf",
    MIME_DOCX,
    MIME_XLSX,
    MIME_PPTX,
    MIME_GDOC,
}

PIPELINE_MARKER_START = "<!-- AUTO-GENERATED-PIPELINE-START -->"
PIPELINE_MARKER_END = "<!-- AUTO-GENERATED-PIPELINE-END -->"

# ── Services ──────────────────────────────────────────────────────────────────

_SERVICES = {}


def get_drive():
    if "drive" in _SERVICES:
        return _SERVICES["drive"]
    with open(CREDS_PATH, "rb") as f:
        creds = pickle.load(f)
    if hasattr(creds, "expired") and creds.expired and creds.refresh_token:
        from google.auth.transport.requests import Request
        creds.refresh(Request())
        with open(CREDS_PATH, "wb") as f:
            pickle.dump(creds, f)
    svc = build("drive", "v3", credentials=creds)
    _SERVICES["drive"] = svc
    _SERVICES["_creds"] = creds
    return svc


def get_docs():
    if "docs" in _SERVICES:
        return _SERVICES["docs"]
    if "_creds" not in _SERVICES:
        get_drive()
    svc = build("docs", "v1", credentials=_SERVICES["_creds"])
    _SERVICES["docs"] = svc
    return svc


# ── Config loading ────────────────────────────────────────────────────────────

def load_drive_docs():
    with open(DRIVE_DOCS_YAML) as f:
        return yaml.safe_load(f)


def get_deal_entry(deal_id):
    cfg = load_drive_docs()
    deal_docs = cfg.get("deal_docs", {})
    if deal_id not in deal_docs:
        die(f"deal '{deal_id}' not found in deal_docs registry")
    return deal_docs[deal_id]


# ── State (extract dedup) ─────────────────────────────────────────────────────

def djb2(s):
    h = 5381
    for c in s:
        h = ((h << 5) + h) + ord(c)
        h &= 0xFFFFFFFF
    return f"{h:08x}"


def load_state():
    if not os.path.exists(EXTRACT_STATE_PATH):
        return {}
    with open(EXTRACT_STATE_PATH) as f:
        return json.load(f)


def save_state(state):
    os.makedirs(os.path.dirname(EXTRACT_STATE_PATH), exist_ok=True)
    with open(EXTRACT_STATE_PATH, "w") as f:
        json.dump(state, f, indent=2)


def deal_state(state, deal_id):
    if deal_id not in state:
        state[deal_id] = {"last_run": None, "processed_files": {}}
    return state[deal_id]


# ── Logging ───────────────────────────────────────────────────────────────────

def die(msg, code=1):
    print(f"ERROR: {msg}", file=sys.stderr)
    sys.exit(code)


def log_error(deal_id, file_id, err):
    os.makedirs(EXTRACT_LOG_DIR, exist_ok=True)
    ts = datetime.now().isoformat()
    with open(ERROR_LOG, "a") as f:
        f.write(f"{ts} deal={deal_id} file={file_id} error={err}\n")


# ── Drive helpers ─────────────────────────────────────────────────────────────

def drive_list_folder(folder_id, recursive=False, _depth=0, _path=""):
    """List files in a Drive folder. When recursive=True, descends into
    subfolders (skipping EXCLUDED_SUBFOLDER_NAMES) up to MAX_FOLDER_DEPTH.
    Each returned file dict gets a synthetic '_path' field showing its
    relative location (e.g. 'Diligence/foo.pdf')."""
    svc = get_drive()
    res = svc.files().list(
        q=f"'{folder_id}' in parents and trashed=false",
        fields="files(id,name,mimeType,modifiedTime,parents)",
        pageSize=200,
    ).execute()
    files = res.get("files", [])
    out = []
    for f in files:
        f["_path"] = (_path + "/" + f["name"]) if _path else f["name"]
        if f["mimeType"] == "application/vnd.google-apps.folder":
            if not recursive or _depth >= MAX_FOLDER_DEPTH:
                out.append(f)  # surface the folder so callers can decide
                continue
            if f["name"] in EXCLUDED_SUBFOLDER_NAMES or f["name"].startswith("_"):
                continue
            out.extend(drive_list_folder(
                f["id"], recursive=True, _depth=_depth + 1, _path=f["_path"]
            ))
        else:
            out.append(f)
    return out


def drive_get_or_create_subfolder(parent_id, name):
    svc = get_drive()
    res = svc.files().list(
        q=f"'{parent_id}' in parents and name='{name}' and "
          f"mimeType='application/vnd.google-apps.folder' and trashed=false",
        fields="files(id,name)",
    ).execute()
    files = res.get("files", [])
    if files:
        return files[0]["id"]
    f = svc.files().create(
        body={"name": name, "mimeType": "application/vnd.google-apps.folder", "parents": [parent_id]},
        fields="id",
    ).execute()
    return f["id"]


def _drive_download_bytes(file_id):
    svc = get_drive()
    request = svc.files().get_media(fileId=file_id)
    buf = io.BytesIO()
    downloader = MediaIoBaseDownload(buf, request)
    done = False
    while not done:
        _, done = downloader.next_chunk()
    return buf.getvalue()


def _docx_to_text(blob_bytes):
    import docx
    doc = docx.Document(io.BytesIO(blob_bytes))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _xlsx_to_text(blob_bytes):
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(blob_bytes), data_only=True, read_only=True)
    out = []
    for sheet in wb.worksheets:
        out.append(f"### Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            if any(c.strip() for c in cells):
                out.append("\t".join(cells))
        out.append("")
    return "\n".join(out)


def _pptx_to_text(blob_bytes):
    from pptx import Presentation
    prs = Presentation(io.BytesIO(blob_bytes))
    out = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"### Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(run.text for run in para.runs).strip()
                    if txt:
                        out.append(txt)
        out.append("")
    return "\n".join(out)


def _pdf_to_text(blob_bytes):
    """Extract text from PDF using pypdf. Page-aware output."""
    import pypdf
    reader = pypdf.PdfReader(io.BytesIO(blob_bytes))
    out = []
    for i, page in enumerate(reader.pages, 1):
        try:
            txt = page.extract_text() or ""
        except Exception:
            txt = ""
        if txt.strip():
            out.append(f"### Page {i}")
            out.append(txt.strip())
            out.append("")
    return "\n".join(out) if out else "[PDF had no extractable text — likely scanned image]"


def drive_read_file_text(file_id):
    """Return file text. Handles native Docs (export), text/plain/markdown,
    .docx/.xlsx/.pptx (parsed with python-docx/openpyxl/python-pptx). PDF
    extraction is not implemented here — returns a marker so the caller
    can decide whether to skip or invoke a separate tool."""
    svc = get_drive()
    meta = svc.files().get(fileId=file_id, fields="id,name,mimeType").execute()
    mime = meta["mimeType"]

    if mime == MIME_GDOC:
        data = svc.files().export(fileId=file_id, mimeType="text/plain").execute()
        return data.decode("utf-8") if isinstance(data, bytes) else data

    if mime == "application/pdf":
        try:
            return _pdf_to_text(_drive_download_bytes(file_id))
        except Exception as e:
            return f"[PDF EXTRACT FAILED — {e} — file_id={file_id}, name={meta['name']}]"

    if mime == MIME_DOCX:
        try:
            return _docx_to_text(_drive_download_bytes(file_id))
        except Exception as e:
            return f"[DOCX EXTRACT FAILED — {e} — file_id={file_id}, name={meta['name']}]"

    if mime == MIME_XLSX:
        try:
            return _xlsx_to_text(_drive_download_bytes(file_id))
        except Exception as e:
            return f"[XLSX EXTRACT FAILED — {e} — file_id={file_id}, name={meta['name']}]"

    if mime == MIME_PPTX:
        try:
            return _pptx_to_text(_drive_download_bytes(file_id))
        except Exception as e:
            return f"[PPTX EXTRACT FAILED — {e} — file_id={file_id}, name={meta['name']}]"

    # text/plain, text/markdown, anything else readable as bytes
    return _drive_download_bytes(file_id).decode("utf-8", errors="replace")


def drive_overwrite_text_file(file_id, text):
    svc = get_drive()
    media = MediaInMemoryUpload(text.encode("utf-8"), mimetype="text/plain")
    svc.files().update(fileId=file_id, media_body=media).execute()


def drive_move(file_id, new_parent_id):
    svc = get_drive()
    f = svc.files().get(fileId=file_id, fields="parents").execute()
    prev = ",".join(f.get("parents", []))
    svc.files().update(
        fileId=file_id, addParents=new_parent_id, removeParents=prev, fields="id,parents"
    ).execute()


def docs_overwrite_native(doc_id, text):
    svc = get_docs()
    doc = svc.documents().get(documentId=doc_id).execute()
    end_index = doc["body"]["content"][-1]["endIndex"] - 1
    requests = []
    if end_index > 1:
        requests.append({"deleteContentRange": {"range": {"startIndex": 1, "endIndex": end_index}}})
    requests.append({"insertText": {"location": {"index": 1}, "text": text}})
    svc.documents().batchUpdate(documentId=doc_id, body={"requests": requests}).execute()


def docs_read_native(doc_id):
    svc = get_docs()
    doc = svc.documents().get(documentId=doc_id).execute()
    out = []
    for elem in doc["body"]["content"]:
        if "paragraph" in elem:
            for pe in elem["paragraph"].get("elements", []):
                if "textRun" in pe:
                    out.append(pe["textRun"]["content"])
    return "".join(out)


# ── Sub-command handlers ──────────────────────────────────────────────────────

def cmd_list_deals(args):
    cfg = load_drive_docs()
    deal_docs = cfg.get("deal_docs", {})
    out = []
    filt = set(args.filter.split(",")) if args.filter else None
    state = load_state()
    for deal_id, entry in deal_docs.items():
        if filt and deal_id not in filt:
            continue
        # Skip non-deal entries (e.g. project_instructions) that lack required deal fields
        if "drive_folder_id" not in entry or "status" not in entry or "master_brief" not in entry:
            continue
        out.append({
            "deal_id": deal_id,
            "drive_folder_id": entry["drive_folder_id"],
            "status_id": entry["status"]["doc_id"],
            "brief_id": entry["master_brief"]["doc_id"],
            "last_run": deal_state(state, deal_id)["last_run"],
        })
    print(json.dumps(out, indent=2))


def cmd_list_new_files(args):
    deal_id = args.deal_id
    entry = get_deal_entry(deal_id)
    folder_id = entry["drive_folder_id"]
    excluded_ids = {entry["status"]["doc_id"], entry["master_brief"]["doc_id"]}

    state = load_state()
    ds = deal_state(state, deal_id)
    last_run = ds.get("last_run")
    processed = ds.get("processed_files", {})

    files = drive_list_folder(folder_id, recursive=True)
    out = []
    for f in files:
        if f["mimeType"] == "application/vnd.google-apps.folder":
            continue  # surfaced excluded folders — drop
        if f["id"] in excluded_ids:
            continue
        if any(f["name"].endswith(s) for s in EXCLUDED_NAME_SUFFIXES):
            continue
        if f["mimeType"] not in PROCESSABLE_MIMES:
            continue
        h = djb2(f["id"])
        prior = processed.get(h)
        if prior:
            # Already touched: success → skip permanently; failed → retry now.
            if prior.get("outcome") == "success":
                continue
            # else: fall through and re-include for retry
        else:
            # Never touched. Use modifiedTime to decide if it's "new since last_run".
            if last_run and f["modifiedTime"] <= last_run:
                continue
        out.append({
            "file_id": f["id"],
            "name": f["name"],
            "path": f.get("_path", f["name"]),
            "mime": f["mimeType"],
            "modified": f["modifiedTime"],
        })
    # Sort by modified time ascending so oldest-first processing is natural.
    out.sort(key=lambda x: x["modified"])
    print(json.dumps(out, indent=2))


def cmd_read_file(args):
    print(drive_read_file_text(args.file_id))


def cmd_write_file_by_id(args):
    """Overwrite a Drive plain-text file by file ID. Reads new content from stdin."""
    text = sys.stdin.read()
    if not text.strip():
        die("empty stdin — refusing to write empty content")
    if args.dry_run:
        print(f"[DRY-RUN] would overwrite file {args.file_id} with {len(text)} chars")
        return
    drive_overwrite_text_file(args.file_id, text)
    print(f"OK wrote {len(text)} chars to {args.file_id}")


_DEAL_DOC_FIELDS = {
    "status":  "status",
    "brief":   "master_brief",
    "lps":     "lps",
    "terms":   "terms",
    "actions": "actions",
}


def cmd_read_deal_doc(args):
    entry = get_deal_entry(args.deal_id)
    field = _DEAL_DOC_FIELDS.get(args.kind)
    if not field:
        die(f"unknown kind '{args.kind}' (expected: {','.join(_DEAL_DOC_FIELDS)})")
    sub = entry.get(field)
    if not sub:
        # File not yet registered for this deal — silent empty, not an error.
        # (e.g. lps/terms/actions may be absent for newly-onboarded deals.)
        print("")
        return
    doc_id = sub["doc_id"]
    svc = get_drive()
    meta = svc.files().get(fileId=doc_id, fields="mimeType").execute()
    if meta["mimeType"] == "application/vnd.google-apps.document":
        print(docs_read_native(doc_id))
    else:
        print(drive_read_file_text(doc_id))


def cmd_write_deal_doc(args):
    entry = get_deal_entry(args.deal_id)
    field = "status" if args.kind == "status" else "master_brief"
    doc_id = entry[field]["doc_id"]
    text = sys.stdin.read()
    if not text.strip():
        die("empty stdin — refusing to write empty doc")
    svc = get_drive()
    meta = svc.files().get(fileId=doc_id, fields="mimeType").execute()
    if args.dry_run:
        print(f"[DRY-RUN] would overwrite {field} doc ({doc_id}) for {args.deal_id} with {len(text)} chars")
        return
    if meta["mimeType"] == "application/vnd.google-apps.document":
        docs_overwrite_native(doc_id, text)
    else:
        drive_overwrite_text_file(doc_id, text)
    print(f"OK wrote {field} for {args.deal_id} ({len(text)} chars)")


def cmd_move_to_ready(args):
    if args.dry_run:
        print(f"[DRY-RUN] would move {args.file_id} into _Ready/ of {args.deal_folder_id}")
        return
    ready_id = drive_get_or_create_subfolder(args.deal_folder_id, "_Ready")
    drive_move(args.file_id, ready_id)
    print(f"OK moved {args.file_id} -> _Ready/ ({ready_id})")


def cmd_mark_processed(args):
    state = load_state()
    ds = deal_state(state, args.deal_id)
    h = djb2(args.file_id)
    ds["processed_files"][h] = {
        "file_id": args.file_id,
        "processed_at": datetime.now().isoformat(),
        "outcome": args.outcome,
    }
    save_state(state)
    print(f"OK marked {args.file_id} {args.outcome}")


def cmd_update_last_run(args):
    state = load_state()
    ds = deal_state(state, args.deal_id)
    ds["last_run"] = datetime.now().isoformat() + "Z"
    save_state(state)
    print(f"OK last_run for {args.deal_id} = {ds['last_run']}")


def cmd_read_log_entries(args):
    """Print new log.json entries for a deal since its last_run.

    log.json is written by cos_capture_pipeline.py (Gmail/Otter/calendar/etc).
    Each entry: {id, date, source, who, what, source_url, source_title, match}.
    /deal-sync reads these as deal-tagged intel feeds — they don't need
    re-extraction; they're already structured."""
    deal_id = args.deal_id
    log_path = os.path.expanduser(f"~/dashboards/data/deals/{deal_id}/log.json")
    if not os.path.exists(log_path):
        print("[]")
        return
    try:
        with open(log_path) as f:
            log = json.load(f)
    except Exception as e:
        die(f"cannot parse {log_path}: {e}")

    entries = log if isinstance(log, list) else log.get("entries", [])

    state = load_state()
    ds = deal_state(state, deal_id)
    last_run = ds.get("last_run")
    captured = ds.get("captured_log_ids", {})

    out = []
    for e in entries:
        if not isinstance(e, dict):
            continue
        eid = str(e.get("id", ""))
        if eid and eid in captured:
            continue
        # Date filter: only entries on/after last_run date
        if last_run and e.get("date"):
            try:
                last_run_date = last_run[:10]  # ISO prefix
                if e["date"] < last_run_date:
                    continue
            except Exception:
                pass
        out.append(e)

    print(json.dumps(out, indent=2))


def cmd_mark_log_captured(args):
    """Record that a log.json entry id has been folded into status/brief.

    Writes two markers (both idempotent):
      1. deal_extract_state.json — captured_log_ids dict for fast lookup
      2. log.json — `folded_at` field stamped on the matching entry, so
         log_compaction.py can archive folded entries deterministically
         (matches the I12 invariant: log.json ≤ 80 KB once entries are
         consumed into status.md + brief.md).
    """
    now_iso = datetime.now().isoformat()

    # 1. External captured_log_ids map (legacy dedup signal)
    state = load_state()
    ds = deal_state(state, args.deal_id)
    if "captured_log_ids" not in ds:
        ds["captured_log_ids"] = {}
    ds["captured_log_ids"][args.entry_id] = now_iso
    save_state(state)

    # 2. Inline `folded_at` on the log.json entry itself
    log_path = os.path.expanduser(
        f"~/dashboards/data/deals/{args.deal_id}/log.json"
    )
    marked_in_log = False
    if os.path.exists(log_path):
        try:
            with open(log_path) as f:
                raw = json.load(f)
        except json.JSONDecodeError:
            raw = None
        if raw is not None:
            # Handle both shapes: {entries: [...]} or [...] at root.
            if isinstance(raw, dict):
                entries = raw.get("entries", [])
            elif isinstance(raw, list):
                entries = raw
            else:
                entries = []
            for e in entries:
                if isinstance(e, dict) and e.get("id") == args.entry_id:
                    # Set folded_at only if not already present (idempotent)
                    if not e.get("folded_at"):
                        e["folded_at"] = now_iso
                        marked_in_log = True
                    break
            if marked_in_log:
                if isinstance(raw, dict):
                    raw["updated_at"] = now_iso
                with open(log_path, "w") as f:
                    json.dump(raw, f, indent=2)

    state_msg = "state+log" if marked_in_log else "state"
    print(f"OK marked {args.entry_id} captured for {args.deal_id} ({state_msg})")


def cmd_append_log_entry(args):
    """Append a new entry to a deal's log.json. Used by intel_capture
    to route ---DEAL-INTEL--- blocks into the canonical deal feed.
    Reads JSON object from stdin."""
    deal_id = args.deal_id
    log_dir = os.path.expanduser(f"~/dashboards/data/deals/{deal_id}")
    log_path = os.path.join(log_dir, "log.json")
    raw = sys.stdin.read()
    try:
        entry = json.loads(raw)
    except Exception as e:
        die(f"invalid JSON: {e}")
    if not isinstance(entry, dict):
        die("entry must be a JSON object")
    if "id" not in entry:
        # Synthesize a stable id from content
        content = (entry.get("title", "") + entry.get("date", "") + entry.get("source", "intel"))
        entry["id"] = djb2(content)
    if "date" not in entry:
        entry["date"] = datetime.now().strftime("%Y-%m-%d")
    if "source" not in entry:
        entry["source"] = "intel"

    os.makedirs(log_dir, exist_ok=True)
    if os.path.exists(log_path):
        try:
            with open(log_path) as f:
                existing = json.load(f)
        except Exception:
            existing = []
    else:
        existing = []
    if not isinstance(existing, list):
        existing = existing.get("entries", []) if isinstance(existing, dict) else []

    # Idempotent: skip if id already present
    if any(e.get("id") == entry["id"] for e in existing):
        print(f"SKIP {entry['id']} already in log")
        return
    existing.append(entry)
    with open(log_path, "w") as f:
        json.dump(existing, f, indent=2)
    print(f"OK appended {entry['id']} to {deal_id}/log.json")


def _resolve_context_folder(entry):
    """Return the deal's _Claude Context folder (where status/brief/entry
    live), falling back to drive_folder_id for legacy/unmigrated tenants."""
    return entry.get("claude_context_folder_id") or entry["drive_folder_id"]


def cmd_read_deal_entry(args):
    """Print the deal's current dashboard_entry.json from Drive, or {} if absent.
    Looks in the deal's _Claude Context/ subfolder; falls back to top-level
    deal folder for tenants that haven't migrated yet."""
    entry = get_deal_entry(args.deal_id)
    folder_id = _resolve_context_folder(entry)
    fname = f"{args.deal_id}_dashboard_entry.json"
    svc = get_drive()
    res = svc.files().list(
        q=f"name='{fname}' and '{folder_id}' in parents and trashed=false",
        fields="files(id,name)",
    ).execute()
    files = res.get("files", [])
    # Backward-compat: if not found in context folder, also check top-level
    if not files and entry.get("claude_context_folder_id"):
        res = svc.files().list(
            q=f"name='{fname}' and '{entry['drive_folder_id']}' in parents and trashed=false",
            fields="files(id,name)",
        ).execute()
        files = res.get("files", [])
    if not files:
        print("{}")
        return
    text = drive_read_file_text(files[0]["id"])
    print(text)


def cmd_write_deal_entry(args):
    """Read JSON from stdin, validate, write to deal's _Claude Context/ folder,
    run sync."""
    import subprocess
    raw = sys.stdin.read()
    if not raw.strip():
        die("empty stdin — refusing to write empty dashboard_entry")
    try:
        entry_obj = json.loads(raw)
    except Exception as e:
        die(f"invalid JSON on stdin: {e}")
    if not isinstance(entry_obj, dict):
        die(f"dashboard_entry must be a JSON object, got {type(entry_obj).__name__}")

    deal_reg = get_deal_entry(args.deal_id)
    folder_id = _resolve_context_folder(deal_reg)
    fname = f"{args.deal_id}_dashboard_entry.json"
    entry_obj["_last_updated_from_session"] = datetime.now().strftime("%Y-%m-%d %H:%M")

    if args.dry_run:
        print(f"[DRY-RUN] would write {fname} ({len(raw)} chars) to folder {folder_id}")
        return

    svc = get_drive()
    res = svc.files().list(
        q=f"name='{fname}' and '{folder_id}' in parents and trashed=false",
        fields="files(id)",
    ).execute()
    files = res.get("files", [])
    body_bytes = json.dumps(entry_obj, indent=2).encode("utf-8")
    media = MediaInMemoryUpload(body_bytes, mimetype="application/json")
    if files:
        svc.files().update(fileId=files[0]["id"], media_body=media).execute()
        action = "updated"
    else:
        svc.files().create(
            body={"name": fname, "parents": [folder_id]},
            media_body=media, fields="id",
        ).execute()
        action = "created"
    print(f"OK {action} {fname} in Drive")

    # Trigger sync to compiled
    sync_script = os.path.expanduser("~/cos-pipeline/tools/sync_deals_from_drive.py")
    try:
        subprocess.run(
            ["/opt/homebrew/bin/python3", sync_script, "--deal-id", args.deal_id],
            timeout=60, check=False,
        )
    except Exception as e:
        print(f"  warning: sync failed (non-fatal): {e}", file=sys.stderr)


_STATUS_MAP = {
    "pending": "open",
    "in progress": "in-progress",
    "in-progress": "in-progress",
    "queued": "open",
    "today": "in-progress",
    "deferred": "open",
    "open": "open",
    "monitor": "open",
    "complete": "closed",
    "closed": "closed",
}


def _map_action_status(s):
    if not s:
        return "open"
    sl = s.lower().strip()
    if sl.startswith("await"):
        return "blocked"
    return _STATUS_MAP.get(sl, "open")


def _is_actions_stub(text):
    """Return True if the file has no real table rows (placeholder only)."""
    return "| # |" not in text


def cmd_write_actions_md(args):
    """Generate actions.md from a dashboard_entry JSON (stdin).

    Reads the `actions` array from the entry, formats it as a 7-column
    markdown table, and writes to:
      1. Local ~/dashboards/data/deals/<deal_id>/actions.md
      2. Drive actions doc via Deal Sync Writer (if deal_sync.yaml configured)

    Skips if the local file already has real table content (hand-maintained).
    """
    deal_id = args.deal_id
    today = datetime.now().strftime("%Y-%m-%d")

    raw = sys.stdin.read()
    if not raw.strip():
        die("empty stdin — pass dashboard_entry JSON on stdin")
    try:
        entry = json.loads(raw)
    except Exception as e:
        die(f"invalid JSON: {e}")

    actions = entry.get("actions", [])
    deal_name = entry.get("name", deal_id)

    local_path = Path(os.path.expanduser(f"~/dashboards/data/deals/{deal_id}/actions.md"))
    existing_text = local_path.read_text() if local_path.exists() else ""

    if not _is_actions_stub(existing_text):
        print(f"SKIP {deal_id}/actions.md — has real content; not overwriting")
        return

    if not actions:
        print(f"SKIP {deal_id} — no actions in entry; leaving stub unchanged")
        return

    lines = [
        f"# {deal_name} — Open Actions",
        "",
        "| # | Action | Owner | Due | Priority | Status | Opened |",
        "|---|---|---|---|---|---|---|",
    ]
    for i, a in enumerate(actions, 1):
        priority = (a.get("priority") or "medium").lower()
        status = _map_action_status(a.get("status") or "")
        due = a.get("due") or "TBD"
        owner = (a.get("owner") or "TBD").replace("|", "/")
        action_text = (a.get("action") or "").replace("|", "/")
        lines.append(f"| {i} | {action_text} | {owner} | {due} | {priority} | {status} | {today} |")

    content = "\n".join(lines) + "\n"

    if args.dry_run:
        print(f"[DRY-RUN] would write {len(actions)} action rows to {local_path}")
        print(content[:600])
        return

    os.makedirs(local_path.parent, exist_ok=True)
    local_path.write_text(content)
    print(f"OK wrote {len(actions)} actions to {local_path}")

    if os.path.exists(DEAL_SYNC_YAML):
        try:
            import requests as _requests
            with open(DEAL_SYNC_YAML) as f:
                ds_cfg = yaml.safe_load(f)
            url = ds_cfg.get("url")
            secret = ds_cfg.get("secret")
            if url and secret:
                deal_reg = get_deal_entry(deal_id)
                actions_sub = deal_reg.get("actions")
                if actions_sub and actions_sub.get("doc_id"):
                    doc_id = actions_sub["doc_id"]
                    resp = _requests.post(url, json={
                        "secret": secret,
                        "fileId": doc_id,
                        "content": content,
                    }, timeout=30)
                    if resp.status_code == 200:
                        print(f"OK synced actions to Drive doc {doc_id}")
                    else:
                        print(f"  warning: Drive sync returned {resp.status_code}", file=sys.stderr)
                else:
                    print(f"  note: no actions doc_id for {deal_id} in drive-docs.yaml — skipped Drive write")
        except Exception as e:
            print(f"  warning: Drive write failed (non-fatal): {e}", file=sys.stderr)
    else:
        print(f"  note: deal_sync.yaml not found — skipped Drive write (deploy Deal Sync Writer to enable)")


DEAL_SYSTEM_DATA_PATH = os.path.expanduser(
    "~/dashboards/data/compiled/deal-system-data.json"
)


def _load_deal_system_data():
    """Load deal-system-data.json. Returns the parsed dict."""
    with open(DEAL_SYSTEM_DATA_PATH) as f:
        return json.load(f)


def cmd_write_jane_brief(args):
    """Write Jane Brief content to a deal's jane_brief Drive Doc via setContent.

    Mirrors the write-deal-doc / write-file-by-id patterns. EP1: never recreates.
    Reads content from --content-file (a local path). If jane_brief_file_id is
    missing from deal-system-data.json, logs a warning and returns 1 so the
    caller can skip gracefully without aborting the whole batch.

    After Drive write, mirrors to local:
        ~/dashboards/data/deals/<deal_id>/jane_brief.md
    using atomic tmp + os.replace to prevent partial writes.
    """
    deal_id = args.deal_id
    content_file = args.content_file

    # Read content from file first (fail early if file missing)
    if not os.path.exists(content_file):
        print(f"ERROR: content-file not found: {content_file}", file=sys.stderr)
        return 1

    with open(content_file, "r", encoding="utf-8") as f:
        content = f.read()

    if not content.strip():
        print(f"ERROR: content-file is empty — refusing to wipe jane_brief for {deal_id}",
              file=sys.stderr)
        return 1

    # Look up deal in deal-system-data.json
    try:
        ds = _load_deal_system_data()
    except Exception as e:
        print(f"ERROR: cannot load deal-system-data.json: {e}", file=sys.stderr)
        return 1

    deal_entry = next(
        (d for d in ds.get("deals", []) if d.get("id") == deal_id), None
    )
    if deal_entry is None:
        print(
            f"WARNING: deal {deal_id!r} not found in deal-system-data.json — skipping",
            file=sys.stderr,
        )
        return 1

    file_id = deal_entry.get("jane_brief_file_id")
    if not file_id:
        print(
            f"WARNING: deal {deal_id!r} has no jane_brief_file_id — "
            f"run bootstrap_jane_drive.py first; skipping",
            file=sys.stderr,
        )
        return 1

    # Dry-run — print intent and exit 0
    if args.dry_run:
        local_path = os.path.expanduser(
            f"~/dashboards/data/deals/{deal_id}/jane_brief.md"
        )
        print(
            f"[dry-run] would setContent on jane_brief Drive Doc {file_id} "
            f"({len(content)} chars) and mirror to {local_path}"
        )
        return 0

    # Drive write — setContent on registered file_id (EP1: never recreate)
    try:
        drive_overwrite_text_file(file_id, content)
    except Exception as e:
        print(f"ERROR: Drive write failed for {deal_id} ({file_id}): {e}", file=sys.stderr)
        return 1

    # Local mirror — atomic write (tmp + os.replace)
    local_path = Path(os.path.expanduser(
        f"~/dashboards/data/deals/{deal_id}/jane_brief.md"
    ))
    local_path.parent.mkdir(parents=True, exist_ok=True)
    tmp_path = local_path.with_suffix(".tmp")
    try:
        tmp_path.write_text(content, encoding="utf-8")
        os.replace(tmp_path, local_path)
    except Exception as e:
        print(f"WARNING: local mirror write failed for {deal_id}: {e}", file=sys.stderr)
        # Non-fatal: Drive write succeeded; local will be pulled on next sync

    print(
        f"OK jane_brief for {deal_id}: Drive {file_id} updated "
        f"({len(content)} chars) + local mirror {local_path}"
    )
    return 0


def cmd_regenerate_pipeline_section(args):
    cfg = load_drive_docs()
    firm_id = cfg["reference_docs"]["firm_context"]["doc_id"]
    new_body = sys.stdin.read()
    if not new_body.strip():
        die("empty stdin — refusing to wipe pipeline section")
    if args.dry_run:
        print(f"[DRY-RUN] would replace pipeline section in firm_context ({firm_id}) with {len(new_body)} chars")
        return
    current = docs_read_native(firm_id)
    start = current.find(PIPELINE_MARKER_START)
    end = current.find(PIPELINE_MARKER_END)
    if start == -1 or end == -1 or end <= start:
        die("pipeline markers not found in firm_context doc")
    end += len(PIPELINE_MARKER_END)
    replaced = (
        current[:start]
        + PIPELINE_MARKER_START
        + "\n"
        + new_body.strip()
        + "\n"
        + PIPELINE_MARKER_END
        + current[end:]
    )
    docs_overwrite_native(firm_id, replaced)
    print(f"OK replaced pipeline section ({len(new_body)} chars)")


# ── CLI ───────────────────────────────────────────────────────────────────────

def build_parser():
    p = argparse.ArgumentParser(description=__doc__.split("\n")[1])
    p.add_argument("--dry-run", action="store_true", help="No Drive writes; print intended action")
    sub = p.add_subparsers(dest="cmd", required=True)

    s = sub.add_parser("list-deals"); s.add_argument("--filter")
    s = sub.add_parser("list-new-files"); s.add_argument("deal_id")
    s = sub.add_parser("read-file"); s.add_argument("file_id")
    s = sub.add_parser("write-file-by-id"); s.add_argument("file_id")
    s = sub.add_parser("read-deal-doc"); s.add_argument("deal_id"); s.add_argument("kind", choices=["status", "brief", "lps", "terms", "actions"])
    s = sub.add_parser("write-deal-doc"); s.add_argument("deal_id"); s.add_argument("kind", choices=["status", "brief"])
    s = sub.add_parser("write-actions-md"); s.add_argument("deal_id")
    s = sub.add_parser("read-deal-entry"); s.add_argument("deal_id")
    s = sub.add_parser("write-deal-entry"); s.add_argument("deal_id")
    s = sub.add_parser("read-log-entries"); s.add_argument("deal_id")
    s = sub.add_parser("mark-log-captured"); s.add_argument("deal_id"); s.add_argument("entry_id")
    s = sub.add_parser("append-log-entry"); s.add_argument("deal_id")
    s = sub.add_parser("move-to-ready"); s.add_argument("file_id"); s.add_argument("deal_folder_id")
    s = sub.add_parser("mark-processed"); s.add_argument("deal_id"); s.add_argument("file_id"); s.add_argument("outcome", choices=["success", "failed"])
    s = sub.add_parser("update-last-run"); s.add_argument("deal_id")
    s = sub.add_parser("regenerate-pipeline-section")

    s = sub.add_parser("write-jane-brief",
                       help="setContent on a deal's jane_brief Drive Doc (EP1 — never recreate)")
    s.add_argument("--deal-id", required=True,
                   help="Deal slug (e.g. the deal id from deal-system-data.json)")
    s.add_argument("--content-file", required=True,
                   help="Path to local file containing the jane_brief markdown content")

    return p


HANDLERS = {
    "list-deals": cmd_list_deals,
    "list-new-files": cmd_list_new_files,
    "read-file": cmd_read_file,
    "write-file-by-id": cmd_write_file_by_id,
    "read-deal-doc": cmd_read_deal_doc,
    "write-deal-doc": cmd_write_deal_doc,
    "write-actions-md": cmd_write_actions_md,
    "read-deal-entry": cmd_read_deal_entry,
    "write-deal-entry": cmd_write_deal_entry,
    "read-log-entries": cmd_read_log_entries,
    "mark-log-captured": cmd_mark_log_captured,
    "append-log-entry": cmd_append_log_entry,
    "move-to-ready": cmd_move_to_ready,
    "mark-processed": cmd_mark_processed,
    "update-last-run": cmd_update_last_run,
    "regenerate-pipeline-section": cmd_regenerate_pipeline_section,
    "write-jane-brief": cmd_write_jane_brief,
}


def main():
    args = build_parser().parse_args()
    HANDLERS[args.cmd](args)


if __name__ == "__main__":
    main()
