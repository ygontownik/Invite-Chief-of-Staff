"""
deck_base.py — Universal deck engine for deal modeling
============================================================
Import this in any deal-specific build_deck_*.py.

Provides:
  - Drawing primitives (add_rect, add_text, add_cell, eyebrow, title, footer)
  - Generic read_model(path, input_names, compute_fn, overrides)
  - run_deck() orchestrator + CLI parser
  - Registry lookup for multi-deal support
"""

import sys
import json
import math
import time
import subprocess
from pathlib import Path
from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Slide dimensions ──────────────────────────────────────────────────────────
W      = Inches(10)
H      = Inches(5.625)
MARGIN = Inches(0.45)

# ── Color palette ─────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1B, 0x2A, 0x4A)
NAVY2  = RGBColor(0x25, 0x38, 0x60)
GOLD   = RGBColor(0x8A, 0x60, 0x00)
AMBER  = RGBColor(0xC4, 0x90, 0x00)
TEAL   = RGBColor(0x1A, 0x4A, 0x5A)
TEALT  = RGBColor(0xE6, 0xF2, 0xF5)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x0A, 0x0A, 0x0A)
CHALK  = RGBColor(0xEB, 0xF0, 0xF8)
LTGOLD = RGBColor(0xFB, 0xF3, 0xDC)
STONE  = RGBColor(0xE2, 0xE8, 0xF2)
OFFWHT = RGBColor(0xFA, 0xFC, 0xFF)
CREAM  = RGBColor(0xF5, 0xF2, 0xED)
MID    = RGBColor(0x5A, 0x6A, 0x7A)
LIGHT  = RGBColor(0x88, 0x99, 0xAA)
WARN   = RGBColor(0xFB, 0xF3, 0xDC)


# ── Drawing primitives ────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_color, border_color=None, border_width=Pt(0.5)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=Pt(10), bold=False, italic=False,
             color=None, align=PP_ALIGN.LEFT,
             v_anchor=None, font_name='Arial',
             fill_color=None, wrap=True, fit=False):
    from pptx.enum.text import MSO_ANCHOR
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    if v_anchor:
        tf.vertical_anchor = v_anchor
    if fill_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = fill_color
    else:
        txBox.fill.background()
    txBox.line.fill.background()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    if color:
        run.font.color.rgb = color
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return txBox


def add_cell(slide, text, x, y, w, h,
             bg=OFFWHT, fg=INK, sz=Pt(9), bold=False, italic=False,
             align=PP_ALIGN.CENTER, border=True, fmt=None):
    border_c = RGBColor(0xCC, 0xCC, 0xCC) if border else None
    add_rect(slide, x, y, w, h, bg, border_c)
    display = _fmt(text, fmt) if fmt else str(text)
    add_text(slide, display, x, y, w, h,
             font_size=sz, bold=bold, color=fg,
             align=align, fill_color=None, wrap=False, fit=False)


def _fmt(val, fmt):
    if not isinstance(val, (int, float)): return str(val)
    if fmt == '$0':   return f"${val:,.0f}M"
    if fmt == '$1':   return f"${val:,.1f}M"
    if fmt == 'pct':  return f"{val:.0f}%"
    if fmt == 'pct1': return f"{val:.1f}%"
    if fmt == '$neg': return f"-${abs(val):,.1f}M"
    return str(val)


def eyebrow(slide, text):
    add_text(slide, text.upper(), MARGIN, Inches(0.22), W - 2*MARGIN, Inches(0.18),
             font_size=Pt(7), color=MID, align=PP_ALIGN.LEFT)


def title(slide, text, y=Inches(0.42), w_override=None):
    w = w_override or W - 2*MARGIN
    add_text(slide, text, MARGIN, y, w, Inches(0.50),
             font_size=Pt(16), bold=True, color=NAVY,
             font_name='Georgia', align=PP_ALIGN.LEFT)


def footer(slide, n, total):
    add_text(slide, f"{n} / {total}",
             W - Inches(0.9), H - Inches(0.25), Inches(0.8), Inches(0.20),
             font_size=Pt(7.5), color=LIGHT, align=PP_ALIGN.RIGHT)
    add_text(slide, "CONFIDENTIAL — TCIP INTERNAL",  # noqa: tenant-leak (TCIP is the product name)
             MARGIN, H - Inches(0.25), Inches(3), Inches(0.20),
             font_size=Pt(7), color=LIGHT, align=PP_ALIGN.LEFT)


def heatmap_color(gp_pct):
    """Return RGBColor for a GP% value: red(25%) → gold(50%) → dark green(65%)."""
    t = (gp_pct - 25) / 40.0
    t = max(0, min(1, t))
    if t <= 0.5:
        u = t / 0.5
        r = int(0x7A + u*(0xC4-0x7A))
        g = int(0x1E + u*(0x90-0x1E))
        b = int(0x1E + u*(0x00-0x1E))
    else:
        u = (t-0.5)/0.5
        r = int(0xC4 + u*(0x16-0xC4))
        g = int(0x90 + u*(0x32-0x90))
        b = int(0x00 + u*(0x28-0x00))
    return RGBColor(r, g, b)


def draw_matrix(slide, M_data, grid, col_labels, row_labels,
                x, y, total_w, row_h,
                header_text, anchor_ri, anchor_ci,
                fmt='$1', header_bg=NAVY,
                row_label_w=Inches(0.72)):
    """Draw a labelled sensitivity matrix."""
    n_cols = len(col_labels)
    col_w = (total_w - row_label_w) / n_cols

    add_rect(slide, x, y, total_w, Inches(0.28), header_bg)
    add_text(slide, header_text,
             x + Inches(0.08), y, total_w - Inches(0.16), Inches(0.28),
             font_size=Pt(7.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_rect(slide, x, y + Inches(0.28), row_label_w, Inches(0.24), NAVY2)
    for ci, hdr in enumerate(col_labels):
        cx = x + row_label_w + ci * col_w
        add_cell(slide, hdr, cx, y + Inches(0.28), col_w, Inches(0.24),
                 bg=NAVY2, fg=WHITE, sz=Pt(6.5), bold=True)
    add_rect(slide, x, y + Inches(0.28), row_label_w, Inches(0.24), NAVY2)
    add_text(slide, "FIT ↓",
             x + Inches(0.04), y + Inches(0.28), row_label_w, Inches(0.24),
             font_size=Pt(7), bold=True, color=WHITE)

    for ri, (row_lbl, row_data) in enumerate(zip(row_labels, grid)):
        ry = y + Inches(0.52) + ri * row_h
        is_fit_anc = (row_lbl in ('$1B', '$1,000M'))
        row_bg = LTGOLD if is_fit_anc else (CHALK if ri % 2 == 0 else OFFWHT)

        add_cell(slide, row_lbl, x, ry, row_label_w, row_h,
                 bg=LTGOLD if is_fit_anc else STONE,
                 fg=GOLD if is_fit_anc else NAVY,
                 sz=Pt(8), bold=is_fit_anc)

        for ci, val in enumerate(row_data):
            cx = x + row_label_w + ci * col_w
            is_anc = (ri == anchor_ri and ci == anchor_ci)
            if fmt == 'heatmap':
                bg = AMBER if is_anc else heatmap_color(val)
                display = f"{val:.1f}%"
                fg_c = WHITE
                extra_txt = "floor" if val <= 25 and not is_anc else None
            else:
                bg = LTGOLD if is_anc else row_bg
                display = _fmt(val, fmt)
                fg_c = GOLD if is_anc else INK
                extra_txt = None

            border_c = AMBER if is_anc else RGBColor(0xCC, 0xCC, 0xCC)
            add_rect(slide, cx, ry, col_w, row_h, bg,
                     border_c, Pt(1.0) if is_anc else Pt(0.5))
            add_text(slide, display, cx, ry, col_w, row_h,
                     font_size=Pt(9) if is_anc else Pt(8.5),
                     bold=is_anc, color=fg_c, align=PP_ALIGN.CENTER)
            if extra_txt:
                add_text(slide, extra_txt, cx, ry + row_h - Inches(0.14),
                         col_w, Inches(0.13),
                         font_size=Pt(5), color=WHITE, align=PP_ALIGN.CENTER)


# ── Generic model reader ──────────────────────────────────────────────────────

def read_model(path: str, input_names: list, compute_fn, overrides: dict = None) -> dict:
    """
    Generic model reader. Deal scripts call this with their own input_names and compute_fn.

    path         : path to .xlsx
    input_names  : list of named range strings to read as plain inputs
    compute_fn   : callable(M, wb) → M — computes all derived values
    overrides    : {NAMED_RANGE: value} applied before compute_fn
    """
    wb = load_workbook(path, data_only=True)
    M  = {}

    for name in input_names:
        if name not in wb.defined_names:
            raise ValueError(f"Named range '{name}' not found in model.")
        ref = wb.defined_names[name].attr_text
        sheet_name, cell_ref = ref.split('!')
        val = wb[sheet_name.strip("'")][cell_ref.replace('$', '')].value
        if val is None:
            raise ValueError(f"Input named range '{name}' is None — model may be corrupted.")
        M[name] = float(val)

    if overrides:
        for k, v in overrides.items():
            M[k] = float(v)

    M = compute_fn(M, wb)
    wb.close()
    return M


# ── Deal registry ─────────────────────────────────────────────────────────────

def _find_registry() -> Path:
    """
    Locate deal_registry.json using standard discovery order:
      1. TCIP_DEAL_REGISTRY env var (explicit override)
      2. ~/cos-pipeline-config-*/config/deal_registry.json (tenant config glob)
      3. Same directory as this file (dev/legacy fallback)
    """
    import os, glob as _glob
    env = os.environ.get('TCIP_DEAL_REGISTRY')
    if env:
        return Path(env)
    candidates = sorted(_glob.glob(
        str(Path.home() / 'cos-pipeline-config-*/config/deal_registry.json')))
    if candidates:
        return Path(candidates[0])
    return Path(__file__).parent / 'deal_registry.json'


def load_registry() -> dict:
    p = _find_registry()
    if not p.exists():
        return {}
    return json.loads(p.read_text())


def registry_entry(deal: str) -> dict:
    reg = load_registry()
    available = [k for k in reg if not k.startswith('_')]
    if deal not in reg:
        raise ValueError(f"Deal '{deal}' not in registry. Available: {available}")
    entry = dict(reg[deal])
    for key in ('model', 'pptx', 'script'):
        if key in entry:
            entry[key] = str(Path(entry[key]).expanduser())
    return entry


def print_new_deal_guide():
    """Print the new-deal startup checklist from registry _meta."""
    reg = load_registry()
    meta = reg.get('_meta', {})
    std = meta.get('standards_doc', {})
    tmpl = meta.get('model_template', '')
    instructions = meta.get('new_deal_instructions', [])

    print("\n── New Deal Startup ─────────────────────────────────────────────")
    if std.get('url'):
        print(f"  Standards doc : {std['url']}")
    if tmpl:
        print(f"  Model template: {tmpl}")
    if instructions:
        print("  Steps:")
        for step in instructions:
            print(f"    {step}")
    print("─────────────────────────────────────────────────────────────────\n")


# ── CLI parser ────────────────────────────────────────────────────────────────

def parse_cli(argv=None):
    """
    Parse standard deck CLI flags.
    Returns (model_path, output_path, updates, fast, only_slides).
    """
    args = list(argv if argv is not None else sys.argv[1:])

    fast = '--fast' in args
    if fast: args.remove('--fast')

    only_slides = None
    if '--slides' in args:
        idx = args.index('--slides')
        only_slides = set(int(x) for x in args[idx+1].split(','))
        args.pop(idx+1); args.pop(idx)
        fast = True

    updates = {}
    if '--update' in args:
        idx = args.index('--update')
        args.pop(idx)
        while idx < len(args) and '=' in args[idx]:
            k, v = args.pop(idx).split('=', 1)
            if k == 'FIT_FEE_BPS':
                k = 'FIT_FEE_PCT'; v = str(float(v) / 10000)
            updates[k] = float(v) if '.' in v else int(v)

    model_path  = args[0] if len(args) > 0 else 'model.xlsx'
    output_path = args[1] if len(args) > 1 else 'output.pptx'

    return model_path, output_path, updates, fast, only_slides


# ── Orchestrator ──────────────────────────────────────────────────────────────

def run_deck(M: dict, output: str, slide_fns: list,
             check_integrity_fn, check_qa_fn,
             fast: bool = False, only_slides: set = None):
    """
    Build a full deck from an already-loaded model dict.

    M                 : populated model dict (from read_*_model)
    output            : path to write .pptx
    slide_fns         : ordered list of slide builder functions
    check_integrity_fn: callable(M) — raises on model errors
    check_qa_fn       : callable(pptx_path, M) — post-build content QA
    fast              : skip visual QA
    only_slides       : set of 1-indexed slide numbers to build (None = all)
    """
    t_start = time.time()

    print("Running integrity checks...")
    check_integrity_fn(M)

    print("Building deck...")
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    total = len(slide_fns)
    for i, fn in enumerate(slide_fns, 1):
        if only_slides and i not in only_slides:
            continue
        fn(prs, M, total)
        if only_slides:
            print(f"  Built slide {i}")

    prs.save(output)
    elapsed = time.time() - t_start
    print(f"✓ Saved: {output}  ({total} slides)  [{elapsed:.1f}s]")

    check_qa_fn(output, M)

    if not fast:
        _convert_to_pdf(output, output.replace('.pptx', '.pdf'))
    else:
        print("(Visual QA skipped — fast mode)")


# ── Environment + recalc helpers ─────────────────────────────────────────────

def _is_claude_env():
    import os
    return os.path.exists('/mnt/skills/public/xlsx/scripts/recalc.py')


def _recalc_model(model_path: str):
    """Force Excel formula recalculation (Mac: xlwings → LibreOffice fallback)."""
    import os
    abs_path = os.path.abspath(model_path)

    if _is_claude_env():
        subprocess.run(
            ['python3', '/mnt/skills/public/xlsx/scripts/recalc.py', model_path, '30'],
            capture_output=True)
        return

    try:
        import xlwings as xw
        app = xw.App(visible=False)
        try:
            wb = app.books.open(abs_path)
            app.calculate()
            wb.save()
            wb.close()
            print("  ✓ Model recalculated via Excel (xlwings)")
            return
        finally:
            app.quit()
    except Exception as e:
        print(f"  ⚠ xlwings recalc failed: {e}")

    lo_paths = ['/Applications/LibreOffice.app/Contents/MacOS/soffice', 'libreoffice', 'soffice']
    for lo in lo_paths:
        try:
            result = subprocess.run(
                [lo, '--headless', '--infilter=Calc MS Excel 2007 XML',
                 '--convert-to', 'xlsx', '--outdir', os.path.dirname(abs_path), model_path],
                capture_output=True, timeout=30)
            if result.returncode == 0:
                return
        except (FileNotFoundError, subprocess.TimeoutExpired):
            continue
    print("  ⚠ Could not recalculate model formulas.")
    print("    Open the model in Excel, press Cmd+Shift+F9, save, then retry.")


def _convert_to_pdf(pptx_path: str, pdf_path: str):
    if _is_claude_env():
        subprocess.run(
            ['python3', '/home/claude/scripts/office/soffice.py',
             '--headless', '--convert-to', 'pdf', pptx_path],
            capture_output=True, cwd='/tmp')
    else:
        lo_paths = ['/Applications/LibreOffice.app/Contents/MacOS/soffice', 'libreoffice', 'soffice']
        for lo in lo_paths:
            try:
                result = subprocess.run(
                    [lo, '--headless', '--convert-to', 'pdf', '--outdir', '/tmp', pptx_path],
                    capture_output=True, timeout=60)
                if result.returncode == 0:
                    return
            except (FileNotFoundError, subprocess.TimeoutExpired):
                continue
        print("  ⚠ PDF conversion failed — LibreOffice not found. Visual QA skipped.")


def update_inputs(model_path: str, changes: dict):
    """Apply input changes directly to the model Excel file (for persistent edits)."""
    wb = load_workbook(model_path)
    applied = []
    for name, value in changes.items():
        if name not in wb.defined_names:
            raise ValueError(f"Named range '{name}' not found in model.")
        ref = wb.defined_names[name].attr_text
        sheet_name, cell_ref = ref.split('!')
        ws = wb[sheet_name.strip("'")]
        old_val = ws[cell_ref.replace('$', '')].value
        ws[cell_ref.replace('$', '')].value = value
        applied.append(f"  {name}: {old_val} → {value}")
    wb.save(model_path)
    print(f"Updated {len(applied)} input(s):")
    for a in applied: print(a)
    return model_path
